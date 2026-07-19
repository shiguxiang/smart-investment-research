"""
PPT 解析器
基于 python-pptx 提取文本框、表格、图片、图表信息
"""

import os
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass, field
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from src.utils.logger import get_logger

logger = get_logger(__name__)


@dataclass
class SlideContent:
    """单张幻灯片内容"""
    slide_num: int
    title: str = ""
    texts: List[str] = field(default_factory=list)
    tables: List[List[List[str]]] = field(default_factory=list)
    images: List[bytes] = field(default_factory=list)
    charts: List[str] = field(default_factory=list)  # 图表描述
    notes: str = ""


@dataclass
class PPTDocument:
    """PPT 文档解析结果"""
    file_path: str
    file_name: str
    total_slides: int
    slides: List[SlideContent] = field(default_factory=list)


class PPTParser:
    """PPT 解析器 - 提取文本、表格、图片、图表"""

    def __init__(self, extract_images: bool = True):
        self.extract_images = extract_images

    def parse(self, file_path: str) -> PPTDocument:
        """
        解析 PPT 文件

        Args:
            file_path: PPT 文件路径

        Returns:
            PPTDocument 包含所有幻灯片的结构化内容
        """
        logger.info(f"开始解析 PPT: {file_path}")

        prs = Presentation(file_path)

        ppt_doc = PPTDocument(
            file_path=file_path,
            file_name=os.path.basename(file_path),
            total_slides=len(prs.slides),
        )

        for slide_idx, slide in enumerate(prs.slides):
            slide_content = self._parse_slide(slide, slide_idx)
            ppt_doc.slides.append(slide_content)

        logger.info(
            f"PPT 解析完成: {ppt_doc.file_name}, "
            f"共 {ppt_doc.total_slides} 张幻灯片"
        )

        return ppt_doc

    def _parse_slide(self, slide, slide_idx: int) -> SlideContent:
        """解析单张幻灯片"""
        content = SlideContent(slide_num=slide_idx + 1)

        for shape in slide.shapes:
            # 1. 标题
            if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
                if shape.has_text_frame:
                    content.title = shape.text_frame.text.strip()

            # 2. 文本框
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    content.texts.append(text)

            # 3. 表格
            elif shape.has_table:
                content.tables.append(self._extract_table(shape.table))

            # 4. 图片
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if self.extract_images:
                    try:
                        image = shape.image
                        content.images.append(image.blob)
                    except Exception as e:
                        logger.warning(
                            f"提取 PPT 图片失败 (slide {slide_idx}): {e}"
                        )

            # 5. 图表
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_desc = self._describe_chart(shape)
                if chart_desc:
                    content.charts.append(chart_desc)

            # 6. SmartArt / 组合图形
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if child.has_text_frame:
                        text = child.text_frame.text.strip()
                        if text:
                            content.texts.append(text)

        # 7. 备注
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            if text_frame:
                content.notes = text_frame.text.strip()

        return content

    @staticmethod
    def _extract_table(table) -> List[List[str]]:
        """提取表格内容"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        return rows

    @staticmethod
    def _describe_chart(shape) -> str:
        """提取图表基本描述"""
        try:
            chart = shape.chart
            chart_type = str(chart.chart_type)
            title = chart.has_title and chart.chart_title.text_frame.text or "无标题"

            categories = []
            series_names = []
            for series in chart.series:
                series_names.append(series.name or "未知系列")

            desc = (
                f"[图表] 类型: {chart_type}, "
                f"标题: {title}, "
                f"数据系列: {', '.join(series_names[:5])}"
            )
            return desc
        except Exception as e:
            logger.debug(f"图表描述提取失败: {e}")
            return "[图表] 无法解析"

    def to_markdown(self, ppt_doc: PPTDocument, include_notes: bool = True) -> str:
        """
        将 PPT 文档转换为 Markdown 格式

        Args:
            ppt_doc: PPTDocument 解析结果
            include_notes: 是否包含备注

        Returns:
            Markdown 格式文本
        """
        lines = [f"# {ppt_doc.file_name}\n"]

        for slide in ppt_doc.slides:
            lines.append(f"## 第 {slide.slide_num} 页")

            if slide.title:
                lines.append(f"### {slide.title}")

            for text in slide.texts:
                lines.append(text)
                lines.append("")

            for table in slide.tables:
                lines.append(self._table_to_markdown(table))
                lines.append("")

            for chart_desc in slide.charts:
                lines.append(f"> {chart_desc}")
                lines.append("")

            if slide.images:
                lines.append(f"> (本页含 {len(slide.images)} 张图片)")

            if include_notes and slide.notes:
                lines.append(f"**备注:** {slide.notes}")

            lines.append("\n---\n")

        return "\n".join(lines)

    @staticmethod
    def _table_to_markdown(table: List[List[str]]) -> str:
        """表格转 Markdown"""
        if not table:
            return ""

        md_lines = []
        # 表头
        md_lines.append("| " + " | ".join(table[0]) + " |")
        md_lines.append("|" + "|".join(["---"] * len(table[0])) + "|")
        # 数据行
        for row in table[1:]:
            # 补齐列数
            padded = row + [""] * (len(table[0]) - len(row))
            md_lines.append("| " + " | ".join(padded[: len(table[0])]) + " |")

        return "\n".join(md_lines)
